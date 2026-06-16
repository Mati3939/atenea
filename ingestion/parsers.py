import re
import unicodedata
from pathlib import Path

SUPPORTED = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def parse(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="ignore")
    return ""


# ── Limpieza de matemática en PDFs (típicamente generados con LaTeX) ────────────

# Glifos de acento sueltos (espaciadores) → combinante equivalente. Solo los que
# usa el español, para no tocar operadores matemáticos (NO incluye ^ ni ~ ASCII).
_ACCENTS = {
    "´": "́", "ˊ": "́",                 # acute  ´
    "˜": "̃",                                       # tilde  ˜  (→ ñ)
    "¨": "̈",                                       # diéresis ¨ (→ ü)
}
_ACCENT_RE = re.compile(
    "([" + "".join(_ACCENTS) + r"])\s?([a-zA-Zıȷ])"
)


def _fix_accents(text: str) -> str:
    """Compone 'contradicci´on' → 'contradicción', 'teor´ıa' → 'teoría'."""
    def repl(m: re.Match) -> str:
        comb = _ACCENTS[m.group(1)]
        base = m.group(2)
        if base == "ı":   # ı (i sin punto, usada por LaTeX)
            base = "i"
        elif base == "ȷ": # ȷ (j sin punto)
            base = "j"
        return unicodedata.normalize("NFC", base + comb)
    return _ACCENT_RE.sub(repl, text)


def _clean_math_text(text: str) -> str:
    text = _fix_accents(text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = re.sub(r" +\n", "\n", text)
    return text


def _parse_pdf(path: Path) -> str:
    import fitz
    doc = fitz.open(path)
    return _clean_math_text("\n".join(page.get_text() for page in doc))


def _parse_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)
